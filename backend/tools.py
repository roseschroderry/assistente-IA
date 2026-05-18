import os
import sys
import shutil
import platform
import subprocess
import json
import time
import requests
import shlex
import re
import base64
import ast
import importlib.util
from io import BytesIO
from typing import Any, Dict, List, Optional
import win32com.client as win32
import psutil # Para monitoramento de sistema
import pyautogui # Para controle direto da interface
from dotenv import load_dotenv
from .app_paths import bundle_dir, core_dir, data_dir, observations_dir, runtime_summary, user_files_dir

# Configuração de segurança do PyAutoGUI
load_dotenv()

PROJECT_ROOT = bundle_dir()
DANGEROUS_TOOL_NAMES = {
    "execute_python_code",
    "run_command",
    "install_missing_dependency",
    "modify_assistant_code",
    "apply_code_change",
}
OPERATOR_TOOL_NAMES = {
    "observe_screen",
    "get_active_window",
    "move_mouse",
    "click_screen",
    "drag_mouse",
    "type_text",
    "press_hotkey",
    "scroll_screen",
}


def _full_access_enabled() -> bool:
    return os.getenv("ELITE_FULL_ACCESS", "").lower() in {"1", "true", "yes", "sim"}


def _dangerous_tools_enabled() -> bool:
    return _full_access_enabled() or os.getenv("ELITE_ENABLE_DANGEROUS_TOOLS", "").lower() in {"1", "true", "yes", "sim"}


def _operator_mode_enabled() -> bool:
    return _full_access_enabled() or os.getenv("ELITE_ENABLE_OPERATOR_MODE", "").lower() in {"1", "true", "yes", "sim"} or _dangerous_tools_enabled()


# Em modo completo, o assistente pode operar a maquina sem a trava de canto do PyAutoGUI.
pyautogui.FAILSAFE = not (
    _full_access_enabled()
    or os.getenv("ELITE_DISABLE_PYAUTOGUI_FAILSAFE", "").lower() in {"1", "true", "yes", "sim"}
)


def _resolve_project_path(file_path: str) -> str:
    abs_path = os.path.abspath(os.path.expandvars(os.path.expanduser(file_path)))
    if _full_access_enabled():
        return abs_path
    project_root = os.path.abspath(PROJECT_ROOT)
    if os.path.commonpath([project_root, abs_path]) != project_root:
        raise PermissionError("Acesso negado a arquivos fora do diretorio do projeto.")
    return abs_path


def _resolve_known_path(path: str) -> str:
    if not path or path == ".":
        return os.getcwd()
    aliases = {
        "desktop": os.path.join(os.path.expanduser("~"), "Desktop"),
        "area de trabalho": os.path.join(os.path.expanduser("~"), "Desktop"),
        "área de trabalho": os.path.join(os.path.expanduser("~"), "Desktop"),
        "documents": os.path.join(os.path.expanduser("~"), "Documents"),
        "documentos": os.path.join(os.path.expanduser("~"), "Documents"),
        "downloads": os.path.join(os.path.expanduser("~"), "Downloads"),
        "projeto": PROJECT_ROOT,
        "project": PROJECT_ROOT,
        "appdata": data_dir(),
        "dados": data_dir(),
        "arquivos": user_files_dir(),
        "files": user_files_dir(),
    }
    normalized = path.lower().strip().replace("\\", "/")
    resolved = aliases.get(path.lower().strip(), path)
    if resolved == path and "/" in normalized:
        first, rest = normalized.split("/", 1)
        if first in aliases and rest:
            original_rest = path.replace("\\", "/").split("/", 1)[1]
            resolved = os.path.join(aliases[first], original_rest)
    return os.path.abspath(os.path.expandvars(os.path.expanduser(resolved)))


def _unique_path(path: str) -> str:
    if not os.path.exists(path):
        return path
    root, ext = os.path.splitext(path)
    for index in range(1, 1000):
        candidate = f"{root}_{index}{ext}"
        if not os.path.exists(candidate):
            return candidate
    return f"{root}_{int(time.time())}{ext}"


def _format_size(size: int) -> str:
    value = float(size)
    for unit in ["B", "KB", "MB", "GB", "TB"]:
        if value < 1024 or unit == "TB":
            return f"{value:.1f} {unit}" if unit != "B" else f"{int(value)} B"
        value /= 1024


def _safe_slug(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9_.-]+", "_", value.strip()).strip("._")
    return slug[:80] or f"snippet_{int(time.time())}"


def _resolve_output_path(filename: str, default_extension: str, path: str = "arquivos", overwrite: bool = False) -> str:
    name = (filename or "").strip() or f"arquivo_{int(time.time())}{default_extension}"
    expanded = os.path.expandvars(os.path.expanduser(name))
    if os.path.isabs(expanded):
        full_path = os.path.abspath(expanded)
    else:
        full_path = os.path.abspath(os.path.join(_resolve_known_path(path or "arquivos"), name))
    if default_extension and not os.path.splitext(full_path)[1]:
        full_path += default_extension
    os.makedirs(os.path.dirname(full_path) or os.getcwd(), exist_ok=True)
    return full_path if overwrite else _unique_path(full_path)


def _clean_sheet_name(value: str, fallback: str = "Dados") -> str:
    cleaned = re.sub(r"[:\\/?*\[\]]+", " ", str(value or fallback)).strip()
    return (cleaned[:31] or fallback)


def _excel_ref_sheet(sheet_name: str) -> str:
    return "'" + sheet_name.replace("'", "''") + "'"


def _infer_table_value(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, (int, float, bool)):
        return value
    text = str(value).strip()
    if not text:
        return ""
    if text.startswith("="):
        return text
    normalized = text.replace("R$", "").replace("$", "").strip()
    if re.fullmatch(r"-?\d+", normalized) and not re.match(r"0\d+", normalized):
        return int(normalized)
    if re.fullmatch(r"-?(\d{1,3}(\.\d{3})+|\d+),\d+", normalized):
        return float(normalized.replace(".", "").replace(",", "."))
    if re.fullmatch(r"-?\d+\.\d+", normalized):
        return float(normalized)
    if re.fullmatch(r"-?\d+%", normalized):
        return float(normalized.rstrip("%")) / 100
    if re.fullmatch(r"-?\d+,\d+%?", normalized):
        number = float(normalized.rstrip("%").replace(",", "."))
        return number / 100 if normalized.endswith("%") else number
    return text


def _parse_table_text(text: str) -> List[List[Any]]:
    raw = (text or "").strip()
    if not raw:
        return [["Item", "Status"], ["Criado pelo Assistente", "OK"]]

    lines = [line.strip(" |") for line in re.split(r"[\r\n;]+", raw) if line.strip(" |")]
    rows = []
    for line in lines:
        if "|" in line:
            parts = [part.strip() for part in line.split("|")]
        elif "\t" in line:
            parts = [part.strip() for part in line.split("\t")]
        elif "," in line:
            parts = [part.strip() for part in line.split(",")]
        else:
            parts = [line.strip()]
        rows.append([_infer_table_value(part) for part in parts if part != ""])

    if not rows:
        return [["Item", "Status"], ["Criado pelo Assistente", "OK"]]
    if len(rows) == 1 and len(rows[0]) == 1:
        return [["Item", "Descricao"], [1, rows[0][0]]]
    return rows


def _coerce_table_data(data: Any) -> List[List[Any]]:
    if isinstance(data, str):
        return _parse_table_text(data)
    if isinstance(data, dict):
        headers = data.get("headers") or data.get("columns") or data.get("colunas")
        rows = data.get("rows") or data.get("data") or data.get("linhas") or []
        if headers:
            return _coerce_table_data([headers] + list(rows))
        if rows:
            return _coerce_table_data(rows)
    if isinstance(data, list) and data and all(isinstance(row, dict) for row in data):
        headers = []
        for row in data:
            for key in row.keys():
                if key not in headers:
                    headers.append(key)
        return [headers] + [[_infer_table_value(row.get(header, "")) for header in headers] for row in data]
    if isinstance(data, list):
        rows = []
        for row in data:
            if isinstance(row, dict):
                rows.extend(_coerce_table_data([row]))
            elif isinstance(row, (list, tuple)):
                rows.append([_infer_table_value(cell) for cell in row])
            else:
                rows.append([_infer_table_value(row)])
        return rows or [["Item", "Status"], ["Criado pelo Assistente", "OK"]]
    return [["Item", "Status"], ["Criado pelo Assistente", "OK"]]


def _apply_excel_polish(ws, table_name: str = "TabelaDados") -> None:
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.worksheet.table import Table, TableStyleInfo
    from openpyxl.utils import get_column_letter

    if ws.max_row < 1 or ws.max_column < 1:
        return

    header_fill = PatternFill("solid", fgColor="202020")
    header_font = Font(bold=True, color="FFFFFF")
    thin = Side(style="thin", color="D9D9D9")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.border = border
    for cell in ws[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.freeze_panes = "A2"

    for column_cells in ws.columns:
        max_length = 8
        for cell in column_cells:
            value = "" if cell.value is None else str(cell.value)
            max_length = max(max_length, min(len(value), 58))
        ws.column_dimensions[get_column_letter(column_cells[0].column)].width = max_length + 3

    header_values = [ws.cell(row=1, column=col).value for col in range(1, ws.max_column + 1)]
    valid_headers = all(isinstance(value, str) and value.strip() for value in header_values)
    if ws.max_row >= 2 and ws.max_column >= 1 and valid_headers:
        ref = f"A1:{get_column_letter(ws.max_column)}{ws.max_row}"
        safe_name = re.sub(r"[^A-Za-z0-9_]", "_", table_name)[:220] or "TabelaDados"
        try:
            for existing in list(ws.tables):
                del ws.tables[existing]
        except Exception:
            pass
        table = Table(displayName=safe_name, ref=ref)
        table.tableStyleInfo = TableStyleInfo(
            name="TableStyleMedium1",
            showFirstColumn=False,
            showLastColumn=False,
            showRowStripes=True,
            showColumnStripes=False,
        )
        ws.add_table(table)


def _add_excel_summary_sheet(wb, source_ws) -> None:
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Font, PatternFill

    if source_ws.max_row < 2 or source_ws.max_column < 1:
        return

    numeric_columns = []
    for col_index in range(1, source_ws.max_column + 1):
        values = [source_ws.cell(row=row, column=col_index).value for row in range(2, source_ws.max_row + 1)]
        numeric_count = sum(isinstance(value, (int, float)) and not isinstance(value, bool) for value in values)
        if numeric_count:
            numeric_columns.append(col_index)

    if "Resumo" in wb.sheetnames:
        del wb["Resumo"]
    ws = wb.create_sheet("Resumo", 0)
    ws["A1"] = "Resumo automatico"
    ws["A1"].font = Font(bold=True, size=16, color="FFFFFF")
    ws["A1"].fill = PatternFill("solid", fgColor="202020")
    ws["A3"] = "Registros"
    ws["B3"] = max(source_ws.max_row - 1, 0)

    row = 5
    ws["A4"] = "Indicador"
    ws["B4"] = "Valor"
    sheet_ref = _excel_ref_sheet(source_ws.title)
    for col_index in numeric_columns:
        header = source_ws.cell(row=1, column=col_index).value or f"Coluna {col_index}"
        col_letter = source_ws.cell(row=1, column=col_index).column_letter
        ws.cell(row=row, column=1, value=f"Total {header}")
        ws.cell(row=row, column=2, value=f"=SUM({sheet_ref}!{col_letter}2:{col_letter}{source_ws.max_row})")
        row += 1
        ws.cell(row=row, column=1, value=f"Media {header}")
        ws.cell(row=row, column=2, value=f"=AVERAGE({sheet_ref}!{col_letter}2:{col_letter}{source_ws.max_row})")
        row += 1

    _apply_excel_polish(ws, "TabelaResumo")
    if numeric_columns and row > 6:
        chart = BarChart()
        chart.title = "Indicadores"
        chart.y_axis.title = "Valor"
        data = Reference(ws, min_col=2, min_row=4, max_row=row - 1)
        cats = Reference(ws, min_col=1, min_row=5, max_row=row - 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 7
        chart.width = 14
        ws.add_chart(chart, "D3")


def _require_operator_mode() -> Optional[str]:
    if _operator_mode_enabled():
        return None
    return "Modo operador desativado. Defina ELITE_ENABLE_OPERATOR_MODE=1 no .env e reinicie o assistente."

def desktop_organize_visual() -> str:
    """Alinha os ícones da área de trabalho à esquerda (simulação via clique direito)."""
    try:
        # Pega o tamanho da tela
        width, height = pyautogui.size()
        
        # Clica com o botão direito no meio da tela (em uma área vazia da área de trabalho)
        # Nota: Isso assume que o dashboard não está cobrindo a tela toda ou que o usuário quer que ele faça isso
        pyautogui.rightClick(width/2, height/2)
        time.sleep(0.5)
        
        # Atalhos de teclado comuns no Windows para organizar ícones:
        # 'v' abre o menu Exibir, depois 'g' ativa "Organizar ícones automaticamente"
        pyautogui.press('v')
        time.sleep(0.2)
        pyautogui.press('g')
        
        return "Comando de organização visual enviado para o Windows."
    except Exception as e:
        return f"Erro ao controlar interface: {str(e)}"

def execute_python_code(code: str) -> str:
    """Executa qualquer codigo Python para dar autonomia total ao assistente."""
    if not _dangerous_tools_enabled():
        return "Ferramenta desativada por seguranca. Defina ELITE_ENABLE_DANGEROUS_TOOLS=1 para habilitar."

    old_stdout = sys.stdout
    try:
        # Captura a saída do código
        import io
        
        redirected_output = sys.stdout = io.StringIO()
        
        # Executa o código no contexto global/local
        exec(code, {"__builtins__": __builtins__}, {})
        
        return f"Código executado. Saída:\n{redirected_output.getvalue()}"
    except Exception as e:
        return f"Erro na execução autônoma: {str(e)}"
    finally:
        sys.stdout = old_stdout

def list_files(directory: str = ".") -> List[str]:
    """Lista arquivos em um diretório com detalhes."""
    try:
        directory = _resolve_known_path(directory)
            
        entries = []
        for name in sorted(os.listdir(directory), key=str.lower):
            path = os.path.join(directory, name)
            if os.path.isdir(path):
                entries.append(f"DIR: {name}")
            else:
                entries.append(f"FILE: {name} ({_format_size(os.path.getsize(path))})")
        return entries
    except Exception as e:
        return [f"Erro ao listar diretório: {str(e)}"]

def get_file_info(path: str) -> str:
    """Mostra metadados basicos de um arquivo ou pasta."""
    try:
        resolved = _resolve_known_path(path)
        if not os.path.exists(resolved):
            return f"Caminho nao encontrado: {resolved}"
        stat = os.stat(resolved)
        kind = "pasta" if os.path.isdir(resolved) else "arquivo"
        size = _format_size(stat.st_size)
        modified = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(stat.st_mtime))
        return f"{kind}: {resolved}\nTamanho: {size}\nModificado: {modified}"
    except Exception as e:
        return f"Erro ao inspecionar arquivo: {str(e)}"

def create_text_file(filename: str, content: str = "", path: str = "arquivos", overwrite: bool = False) -> str:
    """Cria um arquivo de texto em uma pasta conhecida ou caminho completo."""
    try:
        base = _resolve_known_path(path)
        if os.path.isabs(os.path.expandvars(os.path.expanduser(filename))):
            full_path = _resolve_known_path(filename)
        else:
            full_path = os.path.join(base, filename)
        os.makedirs(os.path.dirname(full_path) or os.getcwd(), exist_ok=True)
        if not overwrite:
            full_path = _unique_path(full_path)
        with open(full_path, "w", encoding="utf-8") as f:
            f.write(content or "")
        return f"Arquivo criado em {full_path}."
    except Exception as e:
        return f"Erro ao criar arquivo: {str(e)}"

def append_text_file(file_path: str, content: str) -> str:
    """Acrescenta texto no final de um arquivo."""
    try:
        full_path = _resolve_known_path(file_path)
        os.makedirs(os.path.dirname(full_path) or os.getcwd(), exist_ok=True)
        with open(full_path, "a", encoding="utf-8") as f:
            f.write(content)
            if not content.endswith("\n"):
                f.write("\n")
        return f"Texto adicionado em {full_path}."
    except Exception as e:
        return f"Erro ao adicionar texto: {str(e)}"

def search_files(directory: str = ".", query: str = "", extensions: Optional[List[str]] = None, limit: int = 50) -> List[str]:
    """Procura arquivos por nome em uma pasta."""
    try:
        root = _resolve_known_path(directory)
        needle = (query or "").lower()
        normalized_exts = None
        if extensions:
            normalized_exts = {ext.lower() if ext.startswith(".") else f".{ext.lower()}" for ext in extensions}
        matches = []
        for current_root, dirs, files in os.walk(root):
            dirs[:] = [d for d in dirs if d not in {"__pycache__", ".git", "node_modules", "build", "dist"}]
            for filename in files:
                if needle and needle not in filename.lower():
                    continue
                if normalized_exts and os.path.splitext(filename)[1].lower() not in normalized_exts:
                    continue
                matches.append(os.path.join(current_root, filename))
                if len(matches) >= max(1, min(int(limit), 200)):
                    return matches
        return matches or [f"Nenhum arquivo encontrado em {root}."]
    except Exception as e:
        return [f"Erro ao pesquisar arquivos: {str(e)}"]

def copy_file(source: str, destination: str, overwrite: bool = False) -> str:
    """Copia um arquivo preservando metadados."""
    try:
        src = _resolve_known_path(source)
        dst = _resolve_known_path(destination)
        if os.path.isdir(dst):
            dst = os.path.join(dst, os.path.basename(src))
        os.makedirs(os.path.dirname(dst) or os.getcwd(), exist_ok=True)
        if not overwrite:
            dst = _unique_path(dst)
        shutil.copy2(src, dst)
        return f"Arquivo copiado para {dst}."
    except Exception as e:
        return f"Erro ao copiar arquivo: {str(e)}"

def organize_folder(directory: str) -> str:
    """Organiza arquivos por extensão em pastas (Imagens, Documentos, etc)."""
    try:
        directory = _resolve_known_path(directory)
            
        extensions = {
            'Imagens': ['.jpg', '.jpeg', '.png', '.gif', '.bmp'],
            'Documentos': ['.pdf', '.docx', '.doc', '.txt', '.xlsx', '.pptx'],
            'Executaveis': ['.exe', '.msi'],
            'Compactados': ['.zip', '.rar', '.7z'],
            'Videos': ['.mp4', '.mkv', '.avi']
        }
        
        count = 0
        for file in os.listdir(directory):
            file_path = os.path.join(directory, file)
            if os.path.isfile(file_path):
                ext = os.path.splitext(file)[1].lower()
                for folder, exts in extensions.items():
                    if ext in exts:
                        dest_folder = os.path.join(directory, folder)
                        os.makedirs(dest_folder, exist_ok=True)
                        shutil.move(file_path, os.path.join(dest_folder, file))
                        count += 1
                        break
        return f"Organização concluída! {count} arquivos movidos."
    except Exception as e:
        return f"Erro ao organizar: {str(e)}"

def get_system_stats() -> dict:
    """Retorna estatísticas em tempo real de CPU e RAM."""
    try:
        return {
            "cpu_usage": psutil.cpu_percent(interval=1),
            "ram_usage": psutil.virtual_memory().percent,
            "disk_usage": psutil.disk_usage('/').percent
        }
    except:
        return {"error": "Falha ao obter estatísticas"}

def get_active_window() -> str:
    """Retorna a janela ativa quando o sistema permite consultar essa informacao."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        window = pyautogui.getActiveWindow()
        if not window:
            return "Nao encontrei uma janela ativa."
        return f"Janela ativa: {window.title} | posicao=({window.left},{window.top}) | tamanho={window.width}x{window.height}"
    except Exception as e:
        return f"Erro ao identificar janela ativa: {str(e)}"

def observe_screen() -> str:
    """Captura a tela e registra contexto basico para acompanhar o que esta acontecendo."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        screenshot_dir = observations_dir()
        filename = f"screen_{int(time.time())}.png"
        path = os.path.join(screenshot_dir, filename)
        image = pyautogui.screenshot()
        image.save(path)
        active_window = get_active_window()
        return f"Observacao registrada em {path}. {active_window}"
    except Exception as e:
        return f"Erro ao observar tela: {str(e)}"

def move_mouse(x: int, y: int, duration: float = 0.1) -> str:
    """Move o mouse para uma coordenada absoluta da tela."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        pyautogui.moveTo(x, y, duration=max(0, min(float(duration), 3)))
        return f"Mouse movido para ({x}, {y})."
    except Exception as e:
        return f"Erro ao mover mouse: {str(e)}"

def click_screen(x: int = None, y: int = None, button: str = "left", clicks: int = 1) -> str:
    """Clica na tela na coordenada atual ou informada."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        safe_button = button if button in {"left", "right", "middle"} else "left"
        safe_clicks = max(1, min(int(clicks), 3))
        if x is None or y is None:
            pyautogui.click(button=safe_button, clicks=safe_clicks)
            return f"Clique {safe_button} executado na posicao atual."
        pyautogui.click(x=x, y=y, button=safe_button, clicks=safe_clicks)
        return f"Clique {safe_button} executado em ({x}, {y})."
    except Exception as e:
        return f"Erro ao clicar: {str(e)}"

def drag_mouse(x: int, y: int, duration: float = 0.3, button: str = "left") -> str:
    """Arrasta o mouse ate uma coordenada absoluta da tela."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        safe_button = button if button in {"left", "right", "middle"} else "left"
        pyautogui.dragTo(x, y, duration=max(0, min(float(duration), 5)), button=safe_button)
        return f"Mouse arrastado para ({x}, {y})."
    except Exception as e:
        return f"Erro ao arrastar mouse: {str(e)}"

def type_text(text: str, interval: float = 0.01) -> str:
    """Digita texto na janela ativa."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        pyautogui.write(text, interval=max(0, min(float(interval), 0.2)))
        return f"Texto digitado com {len(text)} caracteres."
    except Exception as e:
        return f"Erro ao digitar texto: {str(e)}"

def press_hotkey(keys: List[str]) -> str:
    """Pressiona uma combinacao de teclas, como ['ctrl', 'c']."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        if not keys:
            return "Nenhuma tecla informada."
        safe_keys = [str(key).lower() for key in keys[:5]]
        pyautogui.hotkey(*safe_keys)
        return f"Atalho executado: {' + '.join(safe_keys)}."
    except Exception as e:
        return f"Erro ao executar atalho: {str(e)}"

def scroll_screen(amount: int) -> str:
    """Rola a janela ativa. Valores positivos sobem, negativos descem."""
    blocked = _require_operator_mode()
    if blocked:
        return blocked
    try:
        safe_amount = max(-2000, min(int(amount), 2000))
        pyautogui.scroll(safe_amount)
        return f"Rolagem executada: {safe_amount}."
    except Exception as e:
        return f"Erro ao rolar tela: {str(e)}"

def create_powerpoint(filename: str, slides_content: List[dict]) -> str:
    """Cria uma apresentação PowerPoint. slides_content: [{'title': '...', 'body': '...'}]"""
    try:
        ppt = win32.gencache.EnsureDispatch('PowerPoint.Application')
        pres = ppt.Presentations.Add()
        
        for slide_data in slides_content:
            slide = pres.Slides.Add(len(pres.Slides) + 1, 1) # 1 = ppLayoutText
            slide.Shapes.Title.TextFrame.TextRange.Text = slide_data.get('title', '')
            slide.Shapes.Placeholders(2).TextFrame.TextRange.Text = slide_data.get('body', '')
            
        path = os.path.abspath(filename)
        pres.SaveAs(path)
        return f"PowerPoint '{filename}' criado com {len(slides_content)} slides."
    except Exception as e:
        return f"Erro ao criar PPT: {str(e)}"

def create_folder(folder_name: str, path: str = ".") -> str:
    """Cria uma nova pasta."""
    path = _resolve_known_path(path)
    full_path = os.path.join(path, folder_name)
    try:
        os.makedirs(full_path, exist_ok=True)
        return f"Pasta '{folder_name}' criada com sucesso em {path}."
    except Exception as e:
        return f"Erro ao criar pasta: {str(e)}"

def move_file(source: str, destination: str) -> str:
    """Move um arquivo ou pasta."""
    try:
        src = _resolve_known_path(source)
        dst = _resolve_known_path(destination)
        if os.path.isdir(dst):
            dst = os.path.join(dst, os.path.basename(src))
        os.makedirs(os.path.dirname(dst) or os.getcwd(), exist_ok=True)
        shutil.move(src, dst)
        return f"Movido de {src} para {dst}."
    except Exception as e:
        return f"Erro ao mover: {str(e)}"

def get_system_info() -> str:
    """Retorna informações do sistema."""
    return f"Sistema: {platform.system()} {platform.release()} | Processador: {platform.processor()}"

def run_command(command: str) -> str:
    """Executa um comando no terminal (cuidado!)."""
    if not _dangerous_tools_enabled():
        return "Ferramenta desativada por seguranca. Defina ELITE_ENABLE_DANGEROUS_TOOLS=1 para habilitar."

    try:
        args = shlex.split(command, posix=False) if platform.system() == "Windows" else shlex.split(command)
        result = subprocess.run(args, shell=False, capture_output=True, text=True, timeout=30)
        return result.stdout if result.returncode == 0 else result.stderr
    except Exception as e:
        return f"Erro ao executar comando: {str(e)}"

# Ferramentas do Office reais
def create_word_doc(filename: str, content: str) -> str:
    """Cria um documento Word."""
    try:
        word = win32.gencache.EnsureDispatch('Word.Application')
        word.Visible = True
        doc = word.Documents.Add()
        range = doc.Range(0, 0)
        range.InsertAfter(content)
        doc.SaveAs(os.path.abspath(filename))
        return f"Documento Word '{filename}' criado e salvo."
    except Exception as e:
        return f"Erro ao criar Word: {str(e)}"

def create_excel_sheet(filename: str, data: List[List[str]], overwrite: bool = False) -> str:
    """Cria uma nova planilha .xlsx sem depender do Excel aberto."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        from openpyxl.utils import get_column_letter

        path = os.path.abspath(os.path.expandvars(os.path.expanduser(filename)))
        if not os.path.splitext(path)[1]:
            path += ".xlsx"
        os.makedirs(os.path.dirname(path) or os.getcwd(), exist_ok=True)
        if not overwrite:
            path = _unique_path(path)

        wb = Workbook()
        ws = wb.active
        ws.title = "Dados"
        for row in data or []:
            ws.append(list(row))
        if ws.max_row >= 1:
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="333333")
            ws.freeze_panes = "A2"
            ws.auto_filter.ref = ws.dimensions
        for column_cells in ws.columns:
            max_length = 8
            for cell in column_cells:
                value = "" if cell.value is None else str(cell.value)
                max_length = max(max_length, min(len(value), 60))
            ws.column_dimensions[get_column_letter(column_cells[0].column)].width = max_length + 2
        wb.save(path)
        return f"Planilha Excel criada em {path}."
    except Exception as fallback_error:
        try:
            excel = win32.gencache.EnsureDispatch('Excel.Application')
            excel.Visible = True
            wb = excel.Workbooks.Add()
            ws = wb.ActiveSheet
            for r_idx, row in enumerate(data or [], 1):
                for c_idx, value in enumerate(row, 1):
                    ws.Cells(r_idx, c_idx).Value = value
            path = os.path.abspath(os.path.expandvars(os.path.expanduser(filename)))
            wb.SaveAs(path)
            return f"Planilha Excel '{filename}' criada com os dados fornecidos."
        except Exception as com_error:
            return f"Erro ao criar Excel: openpyxl={fallback_error}; excel_com={com_error}"

def create_excel_sheet_com(filename: str, data: List[List[str]]) -> str:
    """Cria uma nova planilha usando Microsoft Excel COM quando necessario."""
    try:
        excel = win32.gencache.EnsureDispatch('Excel.Application')
        excel.Visible = True
        wb = excel.Workbooks.Add()
        ws = wb.ActiveSheet
        for r_idx, row in enumerate(data, 1):
            for c_idx, value in enumerate(row, 1):
                ws.Cells(r_idx, c_idx).Value = value
        wb.SaveAs(os.path.abspath(filename))
        return f"Planilha Excel '{filename}' criada com os dados fornecidos."
    except Exception as e:
        return f"Erro ao criar Excel: {str(e)}"


def _coerce_slides(slides_content: Any) -> List[Dict[str, Any]]:
    if isinstance(slides_content, str):
        chunks = [chunk.strip() for chunk in re.split(r"[\r\n;]+", slides_content) if chunk.strip()]
        slides = []
        for chunk in chunks:
            if ":" in chunk:
                title, body = chunk.split(":", 1)
                slides.append({"title": title.strip(), "body": body.strip()})
            else:
                slides.append({"title": chunk[:64], "body": chunk})
        return slides or [{"title": "Apresentacao", "body": slides_content}]
    if isinstance(slides_content, list):
        slides = []
        for index, item in enumerate(slides_content, 1):
            if isinstance(item, dict):
                slides.append({
                    "title": item.get("title") or item.get("titulo") or f"Slide {index}",
                    "body": item.get("body") or item.get("conteudo") or item.get("text") or "",
                    "bullets": item.get("bullets") or item.get("topicos") or [],
                })
            else:
                slides.append({"title": f"Slide {index}", "body": str(item)})
        return slides
    return [{"title": "Apresentacao", "body": "Criada pelo Assistente Elite."}]


def create_powerpoint(filename: str, slides_content: List[dict], path: str = "arquivos", title: str = "", overwrite: bool = False) -> str:
    """Cria uma apresentacao PowerPoint com fallback sem depender do PowerPoint aberto."""
    slides = _coerce_slides(slides_content)
    output_path = _resolve_output_path(filename, ".pptx", path, overwrite)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        if title:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = "Criado pelo Assistente Elite"

        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(slide_data.get("title") or "Slide")
            frame = slide.placeholders[1].text_frame
            frame.clear()
            bullets = slide_data.get("bullets") or []
            body = str(slide_data.get("body") or "")
            items = list(bullets) if bullets else [part.strip() for part in re.split(r"\n+| \u2022 ", body) if part.strip()]
            if not items:
                items = [body or "Conteudo a preencher."]
            for item_index, item in enumerate(items[:8]):
                paragraph = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
                paragraph.text = str(item)
                paragraph.font.size = Pt(22 if item_index == 0 and len(items) == 1 else 18)

            title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
            title_run.font.color.rgb = RGBColor(32, 32, 32)
            title_run.font.bold = True

        prs.save(output_path)
        return f"PowerPoint criado em {output_path} com {len(slides) + (1 if title else 0)} slides."
    except Exception as pptx_error:
        try:
            ppt = win32.gencache.EnsureDispatch('PowerPoint.Application')
            pres = ppt.Presentations.Add()
            for slide_data in slides:
                slide = pres.Slides.Add(len(pres.Slides) + 1, 1)
                slide.Shapes.Title.TextFrame.TextRange.Text = slide_data.get('title', '')
                slide.Shapes.Placeholders(2).TextFrame.TextRange.Text = slide_data.get('body', '')
            pres.SaveAs(output_path)
            return f"PowerPoint criado em {output_path} com {len(slides)} slides."
        except Exception as com_error:
            return f"Erro ao criar PPT: python-pptx={pptx_error}; powerpoint_com={com_error}"


def create_word_doc(filename: str, content: str, title: str = "", path: str = "arquivos", overwrite: bool = False) -> str:
    """Cria um documento Word .docx formatado, com fallback para Word COM."""
    output_path = _resolve_output_path(filename, ".docx", path, overwrite)
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt

        doc = Document()
        styles = doc.styles
        styles["Normal"].font.name = "Segoe UI"
        styles["Normal"].font.size = Pt(11)

        if title:
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for raw_line in (content or "").splitlines() or ["Documento criado pelo Assistente Elite."]:
            line = raw_line.strip()
            if not line:
                doc.add_paragraph()
            elif line.startswith("# "):
                doc.add_heading(line[2:].strip(), level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:].strip(), level=2)
            elif line.startswith(("- ", "* ")):
                doc.add_paragraph(line[2:].strip(), style="List Bullet")
            elif re.match(r"^\d+[.)]\s+", line):
                doc.add_paragraph(re.sub(r"^\d+[.)]\s+", "", line), style="List Number")
            else:
                doc.add_paragraph(line)

        doc.save(output_path)
        return f"Documento Word criado em {output_path}."
    except Exception as docx_error:
        try:
            word = win32.gencache.EnsureDispatch('Word.Application')
            word.Visible = True
            doc = word.Documents.Add()
            doc.Range(0, 0).InsertAfter((title + "\n\n" if title else "") + (content or ""))
            doc.SaveAs(output_path)
            return f"Documento Word criado em {output_path}."
        except Exception as com_error:
            return f"Erro ao criar Word: python-docx={docx_error}; word_com={com_error}"


def create_excel_sheet(
    filename: str,
    data: Any,
    overwrite: bool = False,
    path: str = "arquivos",
    sheet_name: str = "Dados",
    title: str = "",
    include_summary: bool = True,
) -> str:
    """Cria uma planilha .xlsx formatada, com tabela, filtros e resumo automatico."""
    rows = _coerce_table_data(data)
    output_path = _resolve_output_path(filename, ".xlsx", path, overwrite)
    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = _clean_sheet_name(sheet_name)
        if title:
            wb.properties.title = title
        for row in rows:
            ws.append(row)
        _apply_excel_polish(ws, f"Tabela_{_safe_slug(ws.title)}")
        if include_summary:
            _add_excel_summary_sheet(wb, ws)
        wb.save(output_path)
        return f"Planilha Excel criada em {output_path} com {ws.max_row - 1} registros e {ws.max_column} colunas."
    except Exception as fallback_error:
        try:
            excel = win32.gencache.EnsureDispatch('Excel.Application')
            excel.Visible = True
            wb = excel.Workbooks.Add()
            ws = wb.ActiveSheet
            ws.Name = _clean_sheet_name(sheet_name)
            for r_idx, row in enumerate(rows or [], 1):
                for c_idx, value in enumerate(row, 1):
                    ws.Cells(r_idx, c_idx).Value = value
            wb.SaveAs(output_path)
            return f"Planilha Excel criada em {output_path}."
        except Exception as com_error:
            return f"Erro ao criar Excel: openpyxl={fallback_error}; excel_com={com_error}"


def create_excel_workbook(filename: str, sheets: List[dict], path: str = "arquivos", overwrite: bool = False, include_summary: bool = True) -> str:
    """Cria uma pasta de trabalho Excel com varias abas."""
    output_path = _resolve_output_path(filename, ".xlsx", path, overwrite)
    try:
        from openpyxl import Workbook

        wb = Workbook()
        default = wb.active
        wb.remove(default)
        created = []
        for index, sheet in enumerate(sheets or [], 1):
            name = _clean_sheet_name(sheet.get("name") or sheet.get("nome") or f"Aba {index}")
            ws = wb.create_sheet(name)
            rows = _coerce_table_data(sheet.get("data") or sheet.get("rows") or sheet.get("dados") or [])
            for row in rows:
                ws.append(row)
            _apply_excel_polish(ws, f"Tabela_{index}_{_safe_slug(name)}")
            created.append(ws.title)
        if include_summary and created:
            _add_excel_summary_sheet(wb, wb[created[0]])
        wb.save(output_path)
        return f"Workbook Excel criado em {output_path} com {len(created)} abas: {', '.join(created)}."
    except Exception as e:
        return f"Erro ao criar workbook Excel: {str(e)}"


def update_excel_sheet(file_path: str, data: Any, sheet_name: str = "Dados", append: bool = True) -> str:
    """Atualiza uma planilha existente, anexando ou substituindo dados em uma aba."""
    try:
        from openpyxl import load_workbook

        path = _resolve_known_path(file_path)
        rows = _coerce_table_data(data)
        wb = load_workbook(path)
        ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.create_sheet(_clean_sheet_name(sheet_name))
        if not append:
            ws.delete_rows(1, ws.max_row)
            write_rows = rows
        else:
            write_rows = rows[1:] if ws.max_row > 1 and rows else rows
        for row in write_rows:
            ws.append(row)
        _apply_excel_polish(ws, f"Tabela_{_safe_slug(ws.title)}")
        wb.save(path)
        return f"Planilha atualizada em {path}, aba {ws.title}, linhas adicionadas: {len(write_rows)}."
    except Exception as e:
        return f"Erro ao atualizar Excel: {str(e)}"


def web_search(query: str) -> str:
    """Pesquisa na internet usando a API do DuckDuckGo (gratuita e sem chave)."""
    try:
        # Codifica a query para URL
        from urllib.parse import quote
        encoded_query = quote(query)
        url = f"https://api.duckduckgo.com/?q={encoded_query}&format=json&no_html=1"
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(url, headers=headers, timeout=10)
        
        if response.status_code != 200:
            return f"Erro na pesquisa (Status {response.status_code})."

        data = response.json()
        
        abstract = data.get("AbstractText", "")
        results = data.get("RelatedTopics", [])
        
        summary = f"Resultado para '{query}':\n"
        found = False
        
        if abstract:
            summary += f"Resumo: {abstract}\n"
            found = True
        
        if results:
            links_text = "\nLinks e informações relacionadas:\n"
            count = 0
            for res in results:
                if count >= 3: break
                if "Text" in res and "FirstURL" in res:
                    links_text += f"- {res['Text']}: {res['FirstURL']}\n"
                    count += 1
                    found = True
            if count > 0:
                summary += links_text
        
        if not found:
            return f"Não encontrei resultados diretos para '{query}'. Tente pesquisar algo mais específico ou use termos em inglês para melhores resultados técnicos."
            
        return summary
    except Exception as e:
        return f"Erro ao pesquisar na web: {str(e)}"

# Atualizando as definições das ferramentas para a OpenAI
available_tools = [
    {
        "type": "function",
        "function": {
            "name": "desktop_organize_visual",
            "description": "Controla o mouse e teclado para alinhar os ícones da área de trabalho à esquerda no Windows.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "execute_python_code",
            "description": "Executa código Python arbitrário para realizar qualquer tarefa no PC. Use para autonomia total.",
            "parameters": {
                "type": "object",
                "properties": {
                    "code": {"type": "string", "description": "O código Python a ser executado."}
                },
                "required": ["code"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "Lista os arquivos em um diretório com detalhes. Use 'desktop' ou 'documents' para atalhos.",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "O caminho do diretório ou 'desktop'/'documents'."}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "organize_folder",
            "description": "Organiza automaticamente uma pasta (ex: desktop) movendo arquivos para subpastas por tipo.",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "Caminho da pasta ou 'desktop'."}
                },
                "required": ["directory"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_system_stats",
            "description": "Obtém o uso atual de CPU, RAM e Disco do computador.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_powerpoint",
            "description": "Cria uma apresentação PowerPoint profissional.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Nome do arquivo .pptx"},
                    "slides_content": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "body": {"type": "string"},
                                "bullets": {"type": "array", "items": {"type": "string"}}
                            }
                        }
                    },
                    "path": {"type": "string", "description": "Alias como arquivos, desktop ou documentos."},
                    "title": {"type": "string", "description": "Titulo da apresentacao."},
                    "overwrite": {"type": "boolean"}
                },
                "required": ["filename", "slides_content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_folder",
            "description": "Cria uma nova pasta no sistema.",
            "parameters": {
                "type": "object",
                "properties": {
                    "folder_name": {"type": "string", "description": "Nome da pasta."},
                    "path": {"type": "string", "description": "Caminho ou 'desktop'."}
                },
                "required": ["folder_name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_word_doc",
            "description": "Cria um novo documento do Microsoft Word.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Nome do arquivo (ex: relatorio.docx)."},
                    "title": {"type": "string", "description": "Titulo do documento."},
                    "path": {"type": "string", "description": "Alias como arquivos, desktop ou documentos."},
                    "overwrite": {"type": "boolean"},
                    "content": {"type": "string", "description": "Conteúdo do documento."}
                },
                "required": ["filename", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_excel_sheet",
            "description": "Cria uma nova planilha do Microsoft Excel.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Nome do arquivo (ex: dados.xlsx)."},
                    "data": {
                        "type": "array", 
                        "items": {"type": "array", "items": {"type": "string"}},
                        "description": "Matriz de dados para a planilha."
                    },
                    "overwrite": {"type": "boolean", "description": "Quando true, permite sobrescrever o arquivo."},
                    "path": {"type": "string", "description": "Alias como arquivos, desktop ou documentos."},
                    "sheet_name": {"type": "string", "description": "Nome da aba principal."},
                    "title": {"type": "string", "description": "Titulo/metadado da planilha."},
                    "include_summary": {"type": "boolean", "description": "Cria uma aba Resumo com indicadores automaticos."}
                },
                "required": ["filename", "data"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "Pesquisa informações em tempo real na internet.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "O que pesquisar na web."}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "open_application",
            "description": "Abre qualquer programa instalado no PC ou um site.",
            "parameters": {
                "type": "object",
                "properties": {
                    "app_name": {"type": "string", "description": "Nome do executável (ex: notepad.exe) ou URL."}
                },
                "required": ["app_name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "control_system_volume",
            "description": "Aumenta, diminui ou muta o som do computador.",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "enum": ["aumentar", "diminuir", "mudo"]}
                },
                "required": ["action"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "capture_screen",
            "description": "Tira um print da tela inteira para o assistente ver o que está acontecendo.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "install_missing_dependency",
            "description": "Instala automaticamente bibliotecas Python que faltam para executar códigos.",
            "parameters": {
                "type": "object",
                "properties": {
                    "package_name": {"type": "string", "description": "Nome do pacote pip."}
                },
                "required": ["package_name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "learn_user_preference",
            "description": "Memoriza informações importantes sobre o usuário para consultas futuras.",
            "parameters": {
                "type": "object",
                "properties": {
                    "key": {"type": "string", "description": "O que está sendo aprendido (ex: nome, linguagem_favorita)."},
                    "value": {"type": "string", "description": "O valor da informação."}
                },
                "required": ["key", "value"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "recall_user_preferences",
            "description": "Recupera todas as informações memorizadas sobre o usuário.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "modify_assistant_code",
            "description": "Edita ou cria arquivos de código do próprio assistente para auto-evolução.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Caminho do arquivo (ex: backend/tools.py)."},
                    "new_content": {"type": "string", "description": "O novo código completo do arquivo."}
                },
                "required": ["file_path", "new_content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "manage_tasks",
            "description": "Gerencia tarefas e compromissos na agenda do usuário.",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "enum": ["adicionar", "listar", "concluir"]},
                    "task": {"type": "string", "description": "Descrição da tarefa."},
                    "deadline": {"type": "string", "description": "Data ou hora do compromisso."},
                    "task_id": {"type": "integer", "description": "ID da tarefa para marcar como concluída."}
                },
                "required": ["action"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_complex_file",
            "description": "Lê arquivos PDF, Excel ou CSV para análise de dados.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Caminho completo do arquivo."}
                },
                "required": ["file_path"]
            }
        }
    }
]

def manage_tasks(action: str, task: str = None, deadline: str = None, task_id: int = None) -> str:
    """Gerencia a agenda: 'adicionar', 'listar', 'concluir'."""
    from .database import db
    if action == "adicionar":
        return db.add_task(task, deadline)
    elif action == "listar":
        return db.get_tasks()
    elif action == "concluir":
        return db.complete_task(task_id)
    return "Ação de agenda inválida."

def read_complex_file(file_path: str) -> str:
    """Lê o conteúdo de PDFs ou Planilhas Excel para análise."""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            import PyPDF2
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages[:5]: # Limite de 5 páginas para não estourar tokens
                    text += page.extract_text()
                return f"Conteúdo do PDF (limitado):\n{text}"
        elif ext in [".xlsx", ".xls", ".csv"]:
            import pandas as pd
            df = pd.read_excel(file_path) if ext != ".csv" else pd.read_csv(file_path)
            return f"Resumo da Planilha:\n{df.head().to_string()}\nColunas: {list(df.columns)}"
        return f"Formato {ext} não suportado para leitura complexa."
    except ImportError:
        return "Erro: Bibliotecas necessárias (PyPDF2 ou pandas) não instaladas. Use 'install_missing_dependency'."
    except Exception as e:
        return f"Erro ao ler arquivo: {str(e)}"

def read_complex_file(file_path: str) -> str:
    """Le PDFs, Word, PowerPoint, Excel e CSV com resumo mais util para analise."""
    try:
        path = _resolve_known_path(file_path)
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            import PyPDF2
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = "\n".join((page.extract_text() or "") for page in reader.pages[:5])
                return f"PDF: {path}\nPaginas lidas: {min(len(reader.pages), 5)}\nConteudo:\n{text[:5000]}"
        if ext in [".xlsx", ".xls"]:
            import pandas as pd
            book = pd.ExcelFile(path)
            sections = [f"Excel: {path}", f"Abas: {', '.join(book.sheet_names)}"]
            for sheet_name in [name for name in book.sheet_names if name.lower() != "resumo"][:5]:
                df = pd.read_excel(book, sheet_name=sheet_name)
                sections.append(
                    f"\nAba: {sheet_name}\nLinhas: {len(df)} | Colunas: {list(df.columns)}\n"
                    f"Previa:\n{df.head(8).to_string(index=False)}"
                )
                numeric = df.select_dtypes(include="number")
                if not numeric.empty:
                    sections.append(f"Numericos:\n{numeric.describe().round(2).to_string()}")
            return "\n".join(sections)
        if ext == ".csv":
            import pandas as pd
            df = pd.read_csv(path)
            return f"CSV: {path}\nLinhas: {len(df)} | Colunas: {list(df.columns)}\nPrevia:\n{df.head(10).to_string(index=False)}"
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
            return f"Word: {path}\nParagrafos: {len(doc.paragraphs)}\nConteudo:\n{text[:5000]}"
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for index, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                chunks.append(f"Slide {index}: " + " | ".join(texts))
            return f"PowerPoint: {path}\nSlides: {len(prs.slides)}\n" + "\n".join(chunks[:20])
        return f"Formato {ext} nao suportado para leitura complexa."
    except ImportError:
        return "Erro: biblioteca de leitura Office ausente. Rode pip install -r backend\\requirements.txt."
    except Exception as e:
        return f"Erro ao ler arquivo: {str(e)}"


def modify_assistant_code(file_path: str, new_content: str) -> str:
    """Permite que o assistente se auto-modifique ou corrija seus próprios arquivos."""
    if not _dangerous_tools_enabled():
        return "Ferramenta desativada por seguranca. Defina ELITE_ENABLE_DANGEROUS_TOOLS=1 para habilitar."

    try:
        abs_path = _resolve_project_path(file_path)
            
        with open(abs_path, "w", encoding="utf-8") as f:
            f.write(new_content)
        return f"Arquivo {os.path.basename(file_path)} modificado com sucesso. Reinicie para aplicar as mudanças."
    except Exception as e:
        return f"Erro ao modificar código: {str(e)}"

def learn_user_preference(key: str, value: str) -> str:
    """Memoriza algo sobre o usuário (ex: nome, gosto, projeto)."""
    from .database import db
    return db.learn_fact(key, value)

def recall_user_preferences() -> str:
    """Recupera tudo o que o assistente aprendeu sobre o usuário."""
    from .database import db
    facts = db.recall_facts()
    return facts if facts else "Ainda não aprendi nada sobre você."

def open_application(app_name: str) -> str:
    """Abre um aplicativo ou site. Ex: 'notepad', 'chrome', 'https://google.com'."""
    try:
        if app_name.startswith("http"):
            import webbrowser
            webbrowser.open(app_name)
            return f"Abrindo site: {app_name}"
        
        # Tenta abrir como comando de sistema
        os.startfile(app_name)
        return f"Tentando abrir aplicativo: {app_name}"
    except Exception as e:
        return f"Não consegui abrir '{app_name}'. Erro: {str(e)}"

def control_system_volume(action: str) -> str:
    """Controla o volume: 'aumentar', 'diminuir', 'mudo'."""
    try:
        if action == "aumentar":
            for _ in range(5): pyautogui.press("volumeup")
            return "Volume aumentado."
        elif action == "diminuir":
            for _ in range(5): pyautogui.press("volumedown")
            return "Volume diminuído."
        elif action == "mudo":
            pyautogui.press("volumemute")
            return "Volume alternado para mudo."
        return "Ação de volume inválida."
    except Exception as e:
        return f"Erro ao controlar volume: {str(e)}"

def capture_screen() -> str:
    """Tira um print da tela e salva na pasta do assistente."""
    try:
        filename = f"screenshot_{int(time.time())}.png"
        path = os.path.join(observations_dir(), filename)
        pyautogui.screenshot(path)
        return f"Captura de tela salva em: {path}. Eu posso ver sua tela agora!"
    except Exception as e:
        return f"Erro ao capturar tela: {str(e)}"

def capture_screen_payload() -> dict:
    """Captura a tela para exibicao no dashboard."""
    image = pyautogui.screenshot()
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    data_url = "data:image/png;base64," + base64.b64encode(buffer.getvalue()).decode("ascii")
    return {
        "image": data_url,
        "width": image.width,
        "height": image.height,
        "timestamp": time.time(),
    }

def app_diagnostics() -> dict:
    """Retorna diagnostico de prontidao para rodar como app desktop."""
    required_modules = [
        "fastapi",
        "uvicorn",
        "openai",
        "dotenv",
        "win32com",
        "psutil",
        "pyautogui",
        "openpyxl",
        "docx",
        "pptx",
        "xlsxwriter",
        "speech_recognition",
        "pyttsx3",
        "pandas",
        "PyPDF2",
        "pyaudio",
        "PyInstaller",
        "PyQt5",
        "PyQt5.QtWebEngineWidgets",
        "pydantic",
        "requests",
        "pythoncom",
        "pywintypes",
    ]
    modules = {
        module: bool(importlib.util.find_spec(module))
        for module in required_modules
    }
    paths = runtime_summary()
    checks = {
        "frontend_exists": os.path.exists(os.path.join(PROJECT_ROOT, "frontend", "index.html")),
        "data_dir_writable": os.access(data_dir(), os.W_OK),
        "logs_dir_writable": os.access(paths["logs_dir"], os.W_OK),
        "operator_mode": _operator_mode_enabled(),
        "dangerous_tools": _dangerous_tools_enabled(),
        "full_access": _full_access_enabled(),
        "pyautogui_failsafe": pyautogui.FAILSAFE,
    }
    checks["dependencies_ok"] = all(modules.values())
    checks["ready_for_app"] = checks["frontend_exists"] and checks["data_dir_writable"] and checks["dependencies_ok"]
    return {
        "status": "ready" if checks["ready_for_app"] else "needs_attention",
        "paths": paths,
        "checks": checks,
        "modules": modules,
        "browser": browser_status(),
        "tools": len(available_tools) if "available_tools" in globals() else None,
        "system": {
            "platform": platform.platform(),
            "python": sys.version.split()[0],
        },
    }

def install_missing_dependency(package_name: str) -> str:
    """Instala uma biblioteca Python que esteja faltando."""
    if not _dangerous_tools_enabled():
        return "Ferramenta desativada por seguranca. Defina ELITE_ENABLE_DANGEROUS_TOOLS=1 para habilitar."

    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])
        return f"Biblioteca '{package_name}' instalada com sucesso. Tente rodar o comando novamente."
    except Exception as e:
        return f"Erro ao instalar '{package_name}': {str(e)}"

def read_source_code(file_path: str) -> str:
    """Lê o código fonte de um arquivo do assistente."""
    try:
        with open(_resolve_project_path(file_path), 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Erro ao ler arquivo: {str(e)}"

def save_learning_note(title: str, content: str) -> str:
    """Salva uma anotacao de aprendizado para o assistente consultar no projeto."""
    try:
        note_path = os.path.join(core_dir(), "learning_notes.md")
        with open(note_path, "a", encoding="utf-8") as f:
            f.write(f"\n\n## {title}\n\n{content}\n")
        return f"Aprendizado salvo em {note_path}."
    except Exception as e:
        return f"Erro ao salvar aprendizado: {str(e)}"

def integrate_code_snippet(name: str, code: str, notes: str = "") -> str:
    """Guarda um snippet enviado pelo usuario para estudo e reutilizacao futura."""
    try:
        snippets_dir = os.path.join(core_dir(), "snippets")
        os.makedirs(snippets_dir, exist_ok=True)
        slug = _safe_slug(name)
        file_path = os.path.join(snippets_dir, f"{slug}.py")
        with open(file_path, "w", encoding="utf-8") as f:
            if notes:
                f.write(f'"""\n{notes}\n"""\n\n')
            f.write(code)
            if not code.endswith("\n"):
                f.write("\n")
        return f"Snippet integrado em {file_path}."
    except Exception as e:
        return f"Erro ao integrar snippet: {str(e)}"

def apply_code_change(file_path: str, new_code: str) -> str:
    """Aplica uma mudança de código (auto-modificação)."""
    if not _dangerous_tools_enabled():
        return "Ferramenta desativada por seguranca. Defina ELITE_ENABLE_DANGEROUS_TOOLS=1 para habilitar."

    try:
        with open(_resolve_project_path(file_path), 'w', encoding='utf-8') as f:
            f.write(new_code)
        return f"Código em '{file_path}' atualizado com sucesso. Reinicie o sistema para aplicar."
    except Exception as e:
        return f"Erro ao aplicar mudança: {str(e)}"

def scan_machine_index(reason: str = "manual") -> dict:
    """Inicia a varredura inteligente da maquina em segundo plano."""
    from .brain import brain

    return brain.start_background_scan(reason)


def brain_status() -> dict:
    """Mostra o estado atual do cerebro local e do indexador."""
    from .brain import brain

    return brain.status_summary()


def search_brain(query: str, kind: str = "", limit: int = 10) -> list:
    """Busca aplicativos, pastas, arquivos e memorias indexadas pelo cerebro local."""
    from .brain import brain

    return brain.search_items(query=query, kind=kind or None, limit=limit)


def _normalize_open_query(query: str) -> str:
    aliases = {
        "bloco de notas": "notepad",
        "notepad": "notepad",
        "calculadora": "calc",
        "paint": "mspaint",
        "explorador": "explorer",
        "explorer": "explorer",
        "terminal": "wt",
        "powershell": "powershell",
        "cmd": "cmd",
        "prompt": "cmd",
    }
    text = (query or "").strip()
    return aliases.get(text.lower(), text)


def open_by_name(query: str, background: bool = True, kind: str = "") -> str:
    """Abre aplicativo, pasta, planilha ou arquivo pelo nome conhecido no cerebro local."""
    from .brain import brain

    target = _normalize_open_query(query)
    result = brain.open_item(target, background=background, kind=kind or None)
    if result.startswith("Nao encontrei") and target != query:
        return open_application(target)
    if result.startswith("Nao encontrei"):
        fallback = open_application(target)
        if not fallback.startswith("NÃ£o consegui") and not fallback.startswith("Não consegui"):
            return fallback
    return result


def execute_by_name(query: str, background: bool = True) -> str:
    """Executa um aplicativo/script pelo nome, preferindo rodar em segundo plano."""
    return open_by_name(query=query, background=background, kind="")


def remember_brain_note(title: str, content: str, tags: str = "", source_path: str = "") -> str:
    """Salva uma memoria duradoura no cerebro local do assistente."""
    from .brain import brain

    return brain.remember(title=title, content=content, tags=tags, source_path=source_path)


def recall_brain(query: str = "", limit: int = 8) -> str:
    """Recupera memorias salvas no cerebro local."""
    from .brain import brain

    return brain.recall(query=query, limit=limit)


def analyze_and_remember_file(file_path: str, title: str = "", tags: str = "arquivo") -> str:
    """Analisa um arquivo real e salva o resumo no cerebro local."""
    from .brain import brain

    path = _resolve_known_path(file_path)
    summary = read_complex_file(path)
    note_title = title or f"Analise de {os.path.basename(path)}"
    brain.index_path(path, "analise")
    brain.remember(note_title, summary[:6000], tags=tags, source_path=path)
    return f"{summary}\n\nMemoria salva como: {note_title}"


def send_notification(title: str, message: str, channels: Optional[List[str]] = None) -> dict:
    """Envia notificacao por app, Telegram, Pushover, e-mail, Discord, Slack ou webhook WhatsApp."""
    from .notifications import send_notification as send_notification_impl

    return send_notification_impl(title=title, message=message, channels=channels)


def list_notifications(limit: int = 20) -> list:
    """Lista notificacoes recentes do app."""
    from .brain import brain

    return brain.recent_notifications(limit=limit)


def create_flow(name: str, steps: List[Any]) -> str:
    """Salva uma automacao simples para reutilizacao."""
    from .brain import brain

    return brain.save_flow(name=name, steps=steps or [])


def list_flows() -> list:
    """Lista os fluxos de automacao salvos."""
    from .brain import brain

    return brain.list_flows()


def run_flow(name: str) -> str:
    """Executa um fluxo salvo. Aceita passos tool/args, open:, notify: e command:."""
    from .brain import brain

    steps = brain.get_flow(name)
    if steps is None:
        return f"Fluxo '{name}' nao encontrado."
    results = []
    for index, step in enumerate(steps, 1):
        try:
            if isinstance(step, dict):
                tool_name = step.get("tool") or step.get("name")
                args = step.get("args") or {}
                if tool_name not in tool_map:
                    result = f"Ferramenta desconhecida: {tool_name}"
                else:
                    result = tool_map[tool_name](**args)
            else:
                raw = str(step).strip()
                lowered = raw.lower()
                if lowered.startswith("open:"):
                    result = open_by_name(raw.split(":", 1)[1].strip())
                elif lowered.startswith("notify:"):
                    result = send_notification("Assistente Elite", raw.split(":", 1)[1].strip(), ["app"])
                elif lowered.startswith("command:"):
                    result = run_command(raw.split(":", 1)[1].strip())
                else:
                    result = "Passo ignorado: use open:, notify:, command: ou objeto tool/args."
        except Exception as exc:
            result = f"Erro: {exc}"
        results.append(f"{index}. {result}")
    return "\n".join(results)

def browser_status() -> dict:
    """Mostra status da camada Computer Use / Browserbase / Stagehand."""
    from .computer_browser import computer_browser

    return computer_browser.status()


def browser_start_session(goal: str = "", start_url: str = "", mode: str = "read", provider: str = "auto") -> dict:
    """Cria uma sessao de navegador operacional com politica de seguranca."""
    from .computer_browser import computer_browser

    return computer_browser.create_session(goal=goal, start_url=start_url, mode=mode, provider=provider)


def browser_fetch_page(url: str, limit: int = 6000) -> dict:
    """Le uma pagina via navegador/fetch seguro e devolve texto limpo para analise."""
    from .computer_browser import computer_browser

    return computer_browser.fetch_page(url=url, limit=limit)


def browser_run_instruction(instruction: str, url: str = "", mode: str = "read", session_id: str = "") -> dict:
    """Executa ou prepara uma instrucao de navegador. Acoes sensiveis viram pedido de aprovacao."""
    from .computer_browser import computer_browser

    return computer_browser.run_instruction(instruction=instruction, url=url, mode=mode, session_id=session_id or None)


def browser_pending_approvals(limit: int = 20) -> list:
    """Lista acoes de navegador aguardando aprovacao humana."""
    from .computer_browser import computer_browser

    return computer_browser.pending_approvals(limit=limit)


def browser_approve_action(approval_id: str, approved: bool, note: str = "") -> dict:
    """Aprova ou rejeita uma acao pendente do navegador operacional."""
    from .computer_browser import computer_browser

    return computer_browser.decide_approval(approval_id=approval_id, approved=approved, note=note)

# Adicionando ao available_tools
available_tools.extend([
    {
        "type": "function",
        "function": {
            "name": "read_source_code",
            "description": "Lê o código fonte de um arquivo do assistente para análise.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Caminho do arquivo .py ou .html"}
                },
                "required": ["file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "apply_code_change",
            "description": "Modifica o próprio código fonte do assistente (auto-evolução).",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Caminho do arquivo a ser modificado."},
                    "new_code": {"type": "string", "description": "O novo código completo do arquivo."}
                },
                "required": ["file_path", "new_code"]
            }
        }
    }
])

available_tools.extend([
    {
        "type": "function",
        "function": {
            "name": "observe_screen",
            "description": "Captura a tela atual e informa a janela ativa para acompanhar o que o usuario esta fazendo.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_active_window",
            "description": "Retorna o titulo, posicao e tamanho da janela ativa.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "move_mouse",
            "description": "Move o mouse para coordenadas absolutas da tela.",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer"},
                    "y": {"type": "integer"},
                    "duration": {"type": "number"}
                },
                "required": ["x", "y"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "click_screen",
            "description": "Clica na coordenada informada ou na posicao atual do mouse.",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer"},
                    "y": {"type": "integer"},
                    "button": {"type": "string", "enum": ["left", "right", "middle"]},
                    "clicks": {"type": "integer"}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "drag_mouse",
            "description": "Arrasta o mouse ate coordenadas absolutas da tela.",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer"},
                    "y": {"type": "integer"},
                    "duration": {"type": "number"},
                    "button": {"type": "string", "enum": ["left", "right", "middle"]}
                },
                "required": ["x", "y"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "type_text",
            "description": "Digita texto na janela ativa.",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string"},
                    "interval": {"type": "number"}
                },
                "required": ["text"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "press_hotkey",
            "description": "Executa combinacoes de teclado, como ctrl+c, win+d ou alt+tab.",
            "parameters": {
                "type": "object",
                "properties": {
                    "keys": {"type": "array", "items": {"type": "string"}}
                },
                "required": ["keys"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "scroll_screen",
            "description": "Rola a janela ativa. Positivo sobe; negativo desce.",
            "parameters": {
                "type": "object",
                "properties": {
                    "amount": {"type": "integer"}
                },
                "required": ["amount"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "save_learning_note",
            "description": "Salva uma nota de aprendizado persistente dentro do projeto.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "content": {"type": "string"}
                },
                "required": ["title", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "integrate_code_snippet",
            "description": "Guarda um codigo enviado pelo usuario como snippet reutilizavel no cerebro local do assistente.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "code": {"type": "string"},
                    "notes": {"type": "string"}
                },
                "required": ["name", "code"]
            }
        }
    }
])

available_tools.extend([
    {
        "type": "function",
        "function": {
            "name": "move_file",
            "description": "Move um arquivo ou pasta para outro caminho, aceitando aliases como desktop, documentos e arquivos.",
            "parameters": {
                "type": "object",
                "properties": {
                    "source": {"type": "string"},
                    "destination": {"type": "string"}
                },
                "required": ["source", "destination"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_system_info",
            "description": "Mostra informacoes basicas do sistema operacional e processador.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": "Executa um comando local no terminal. Use apenas para tarefas solicitadas pelo dono.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "Comando a executar."}
                },
                "required": ["command"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_excel_workbook",
            "description": "Cria um Excel com varias abas, tabelas formatadas e resumo automatico.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "sheets": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string"},
                                "data": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}}
                            }
                        }
                    },
                    "path": {"type": "string"},
                    "overwrite": {"type": "boolean"},
                    "include_summary": {"type": "boolean"}
                },
                "required": ["filename", "sheets"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "update_excel_sheet",
            "description": "Atualiza uma planilha existente anexando ou substituindo linhas em uma aba.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string"},
                    "data": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}},
                    "sheet_name": {"type": "string"},
                    "append": {"type": "boolean"}
                },
                "required": ["file_path", "data"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_file_info",
            "description": "Mostra metadados de um arquivo ou pasta, como caminho, tamanho e data de modificacao.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Caminho, alias como desktop/documentos, ou arquivo."}
                },
                "required": ["path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_text_file",
            "description": "Cria um arquivo de texto com conteudo informado em uma pasta conhecida ou caminho completo.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "content": {"type": "string"},
                    "path": {"type": "string", "description": "Alias como arquivos, desktop, documentos, projeto ou caminho completo."},
                    "overwrite": {"type": "boolean"}
                },
                "required": ["filename"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "append_text_file",
            "description": "Acrescenta texto no final de um arquivo.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string"},
                    "content": {"type": "string"}
                },
                "required": ["file_path", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_files",
            "description": "Pesquisa arquivos por nome e extensao dentro de uma pasta.",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string"},
                    "query": {"type": "string"},
                    "extensions": {"type": "array", "items": {"type": "string"}},
                    "limit": {"type": "integer"}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "copy_file",
            "description": "Copia um arquivo para outro caminho preservando metadados.",
            "parameters": {
                "type": "object",
                "properties": {
                    "source": {"type": "string"},
                    "destination": {"type": "string"},
                    "overwrite": {"type": "boolean"}
                },
                "required": ["source", "destination"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "app_diagnostics",
            "description": "Executa diagnostico de prontidao do assistente para rodar como app desktop.",
            "parameters": {"type": "object", "properties": {}}
        }
    }
])

available_tools.extend([
    {
        "type": "function",
        "function": {
            "name": "scan_machine_index",
            "description": "Inicia ou acompanha a varredura da maquina para indexar apps, pastas, planilhas e arquivos por nome.",
            "parameters": {
                "type": "object",
                "properties": {
                    "reason": {"type": "string", "description": "Motivo da varredura, como startup ou manual."}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "brain_status",
            "description": "Mostra estado do cerebro local: varredura, itens indexados, memorias e fluxos.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_brain",
            "description": "Busca no cerebro local por aplicativo, pasta, arquivo, planilha, documento ou codigo pelo nome.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"},
                    "kind": {"type": "string", "description": "Opcional: app, folder, spreadsheet, document, presentation, pdf, code, text."},
                    "limit": {"type": "integer"}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "open_by_name",
            "description": "Abre app, pasta, planilha ou arquivo pelo nome usando o indice do cerebro local.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"},
                    "background": {"type": "boolean"},
                    "kind": {"type": "string"}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "execute_by_name",
            "description": "Executa aplicativo ou script pelo nome, preferindo segundo plano.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"},
                    "background": {"type": "boolean"}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "remember_brain_note",
            "description": "Salva uma memoria permanente no cerebro local.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "content": {"type": "string"},
                    "tags": {"type": "string"},
                    "source_path": {"type": "string"}
                },
                "required": ["title", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "recall_brain",
            "description": "Recupera memorias salvas no cerebro local por texto livre.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"},
                    "limit": {"type": "integer"}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "analyze_and_remember_file",
            "description": "Le e resume arquivo, planilha, PDF, Word, PowerPoint, codigo ou CSV e salva no cerebro local.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string"},
                    "title": {"type": "string"},
                    "tags": {"type": "string"}
                },
                "required": ["file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_notification",
            "description": "Envia notificacao por app, Telegram, Pushover, email, Discord, Slack ou webhook WhatsApp configurado.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "message": {"type": "string"},
                    "channels": {"type": "array", "items": {"type": "string"}}
                },
                "required": ["title", "message"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_notifications",
            "description": "Lista notificacoes recentes do app.",
            "parameters": {
                "type": "object",
                "properties": {"limit": {"type": "integer"}}
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_flow",
            "description": "Salva uma automacao de passos reutilizaveis.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "steps": {
                        "type": "array",
                        "items": {"type": "string"}
                    }
                },
                "required": ["name", "steps"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_flows",
            "description": "Lista fluxos de automacao salvos.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "run_flow",
            "description": "Executa um fluxo de automacao salvo.",
            "parameters": {
                "type": "object",
                "properties": {"name": {"type": "string"}},
                "required": ["name"]
            }
        }
    }
])

available_tools.extend([
    {
        "type": "function",
        "function": {
            "name": "browser_status",
            "description": "Mostra status da camada de navegador operacional, Browserbase, Stagehand, allowlist e aprovacoes pendentes.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "browser_start_session",
            "description": "Cria uma sessao isolada de navegador operacional para uma meta ou site.",
            "parameters": {
                "type": "object",
                "properties": {
                    "goal": {"type": "string"},
                    "start_url": {"type": "string"},
                    "mode": {"type": "string", "enum": ["read", "prepare", "approval"]},
                    "provider": {"type": "string", "enum": ["auto", "local-fetch", "browserbase-stagehand"]}
                }
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "browser_fetch_page",
            "description": "Le pagina web, extrai texto limpo e retorna resumo bruto para pesquisa/analise sem clicar.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string"},
                    "limit": {"type": "integer"}
                },
                "required": ["url"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "browser_run_instruction",
            "description": "Opera ou prepara uma tarefa no navegador por linguagem natural. Acoes de alto impacto exigem aprovacao humana.",
            "parameters": {
                "type": "object",
                "properties": {
                    "instruction": {"type": "string"},
                    "url": {"type": "string"},
                    "mode": {"type": "string", "enum": ["read", "prepare", "approval"]},
                    "session_id": {"type": "string"}
                },
                "required": ["instruction"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "browser_pending_approvals",
            "description": "Lista acoes do navegador operacional aguardando aprovacao humana.",
            "parameters": {
                "type": "object",
                "properties": {"limit": {"type": "integer"}}
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "browser_approve_action",
            "description": "Aprova ou rejeita uma acao pendente do navegador operacional pelo ID.",
            "parameters": {
                "type": "object",
                "properties": {
                    "approval_id": {"type": "string"},
                    "approved": {"type": "boolean"},
                    "note": {"type": "string"}
                },
                "required": ["approval_id", "approved"]
            }
        }
    }
])

if not _dangerous_tools_enabled():
    available_tools = [
        tool for tool in available_tools
        if tool.get("function", {}).get("name") not in DANGEROUS_TOOL_NAMES
    ]

if not _operator_mode_enabled():
    available_tools = [
        tool for tool in available_tools
        if tool.get("function", {}).get("name") not in OPERATOR_TOOL_NAMES
    ]

# Dicionario de mapeamento criado depois de todas as funcoes.
tool_map = {
    "desktop_organize_visual": desktop_organize_visual,
    "execute_python_code": execute_python_code,
    "list_files": list_files,
    "get_file_info": get_file_info,
    "create_text_file": create_text_file,
    "append_text_file": append_text_file,
    "search_files": search_files,
    "copy_file": copy_file,
    "organize_folder": organize_folder,
    "get_system_stats": get_system_stats,
    "app_diagnostics": app_diagnostics,
    "create_powerpoint": create_powerpoint,
    "create_folder": create_folder,
    "move_file": move_file,
    "get_system_info": get_system_info,
    "run_command": run_command,
    "create_word_doc": create_word_doc,
    "create_excel_sheet": create_excel_sheet,
    "create_excel_workbook": create_excel_workbook,
    "update_excel_sheet": update_excel_sheet,
    "web_search": web_search,
    "open_application": open_application,
    "control_system_volume": control_system_volume,
    "capture_screen": capture_screen,
    "observe_screen": observe_screen,
    "get_active_window": get_active_window,
    "move_mouse": move_mouse,
    "click_screen": click_screen,
    "drag_mouse": drag_mouse,
    "type_text": type_text,
    "press_hotkey": press_hotkey,
    "scroll_screen": scroll_screen,
    "install_missing_dependency": install_missing_dependency,
    "learn_user_preference": learn_user_preference,
    "recall_user_preferences": recall_user_preferences,
    "modify_assistant_code": modify_assistant_code,
    "manage_tasks": manage_tasks,
    "read_complex_file": read_complex_file,
    "read_source_code": read_source_code,
    "apply_code_change": apply_code_change,
    "save_learning_note": save_learning_note,
    "integrate_code_snippet": integrate_code_snippet,
    "scan_machine_index": scan_machine_index,
    "brain_status": brain_status,
    "search_brain": search_brain,
    "open_by_name": open_by_name,
    "execute_by_name": execute_by_name,
    "remember_brain_note": remember_brain_note,
    "recall_brain": recall_brain,
    "analyze_and_remember_file": analyze_and_remember_file,
    "send_notification": send_notification,
    "list_notifications": list_notifications,
    "create_flow": create_flow,
    "list_flows": list_flows,
    "run_flow": run_flow,
    "browser_status": browser_status,
    "browser_start_session": browser_start_session,
    "browser_fetch_page": browser_fetch_page,
    "browser_run_instruction": browser_run_instruction,
    "browser_pending_approvals": browser_pending_approvals,
    "browser_approve_action": browser_approve_action,
}
